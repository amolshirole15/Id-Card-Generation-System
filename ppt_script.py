from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE  # Import MSO_SHAPE

# Initialize presentation
presentation = Presentation()

# Function to add a slide with title and content
def add_slide_with_content(title, content):
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title_box = slide.shapes.title
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = content
    p.font.size = Pt(20)
    p.font.bold = False
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Function to add a blank slide with a placeholder for screenshots
def add_screenshot_slide(title):
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title_box = slide.shapes.title
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(36)
    title_box.text_frame.paragraphs[0].font.bold = True

    # Add a placeholder rectangle for screenshots
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3.5)
    placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    placeholder.text = "Screenshot Placeholder"
    placeholder.text_frame.paragraphs[0].font.size = Pt(20)
    placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add slides
add_slide_with_content(
    "ID Card Management System",
    "A Comprehensive Solution for Employees and Students\n\nBy: Your Name\nYour College Name\nPresentation Date"
)

add_slide_with_content(
    "Project Overview",
    "- A Django-based application to manage employee/student data.\n"
    "- Create and print ID cards with unique QR codes.\n"
    "- Efficient search functionality with QR code scanning.\n"
    "- Generate high-quality PDF versions of ID cards.\n"
    "- Supports multiple organizations with custom templates."
)

add_slide_with_content(
    "Key Features",
    "- **User Registration & Profile Management:**\n  - Create or update user profiles.\n  - Assign unique organization-specific codes.\n"
    "- **ID Card Templates:**\n  - Multiple templates for ID card customization.\n"
    "- **Employee/Student Management:**\n  - Add, edit, delete, and view employee/student details."
)

add_slide_with_content(
    "ID Card Functionality",
    "- **ID Card Generation:**\n  - Automatically generates cards with custom templates.\n  - Includes employee/student details and unique QR codes.\n"
    "- **Preview and Print:**\n  - Preview ID cards before printing.\n  - High-quality downloadable PDFs."
)

add_slide_with_content(
    "Search & QR Code Scanning",
    "- **Search Functionality:**\n  - Search employees/students using names, codes, or departments.\n"
    "- **QR Code Scanning:**\n  - Instant retrieval of employee/student details by scanning QR codes.\n  - Enhances efficiency in large organizations."
)

add_slide_with_content(
    "Technical Implementation",
    "- **Backend:** Django 4.0 with a modular structure.\n"
    "- **Database:** SQLite/PostgreSQL for robust data storage.\n"
    "- **Libraries:**\n  - Pillow: Image manipulation.\n  - django-qr-code: QR code generation.\n  - python-pptx: Presentation creation.\n"
    "- **Frontend:** HTML, CSS, Bootstrap for user-friendly design."
)

add_slide_with_content(
    "System Workflow",
    "- **Admin Panel:**\n  - Manage users, templates, and profiles.\n"
    "- **Employee/Student Module:**\n  - Add, update, or delete records.\n  - Assign unique employee codes linked to profiles.\n"
    "- **ID Card Module:**\n  - Generate, preview, and print ID cards."
)

add_screenshot_slide("Screenshots of the Project")

add_slide_with_content(
    "Benefits of the System",
    "- Streamlined management of employee/student data.\n"
    "- Time-saving with QR code scanning and search functionality.\n"
    "- Easy customization of ID cards for various organizations.\n"
    "- Secure and efficient data handling."
)

add_slide_with_content(
    "Thank You",
    "Questions?\n\nContact Information:\nEmail: your.email@example.com\nPhone: +1234567890\n\nPresentation By: Your Name"
)

# Save presentation
presentation.save("Enhanced_ID_Card_Management_System_Presentation.pptx")
print("Presentation created: Enhanced_ID_Card_Management_System_Presentation.pptx")
